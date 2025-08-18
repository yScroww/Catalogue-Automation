import os
import sys
import shutil
import math
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
from io import BytesIO

# Importações corrigidas para serem relativas
from .utils.excel import load_product_data, load_image_links
from .utils.images import find_image_for_sku, download_image
from .utils.logger import get_logger

logger = get_logger(__name__)

def add_textbox(slide, left, top, width, height, text, font_size, bold=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    font = run.font
    font.size = font_size
    font.bold = bold
    font.name = "Calibri" 
    p.alignment = align
    return box

# A função agora recebe a marca (brand) como argumento
def add_product_tile(slide, left, top, box_w, box_h, image_path, nome, sku, brand, layout_config):
    titulo_h = layout_config['PRODUCT_NAME_BOX_HEIGHT']
    sku_h = layout_config['SKU_BOX_HEIGHT']
    brand_h = layout_config['BRAND_BOX_HEIGHT']
    padding = layout_config['IMAGE_PADDING']

    img_top = top + padding
    img_left = left + padding
    img_w = box_w - 2 * padding
    img_h = box_h - titulo_h - sku_h - brand_h - 2 * padding

    # Verifica se o caminho da imagem existe antes de tentar inserir
    if image_path and os.path.exists(image_path):
        try:
            with Image.open(image_path) as im:
                # CORREÇÃO: Salva a imagem em um buffer de memória como JPEG
                image_stream = BytesIO()
                im.save(image_stream, format='JPEG', quality=90)  # Qualidade 90%
                image_stream.seek(0)
                
                # Calcule o tamanho da imagem e adicione-a a partir do buffer
                iw, ih = im.size
                scale = min(img_w / iw, img_h / ih)
                new_w = Inches(iw * scale)
                new_h = Inches(ih * scale)
                slide.shapes.add_picture(
                    image_stream,
                    img_left + (img_w - new_w) / 2,
                    img_top + (img_h - new_h) / 2,
                    width=new_w, height=new_h
                )
        except Exception as e:
            logger.warning(f"Falha ao inserir imagem '{image_path}': {e}")
            add_textbox(slide, img_left, img_top, img_w, img_h, "Imagem inválida", font_size=layout_config['SKU_FONT_SIZE'])
    else:
        add_textbox(slide, img_left, img_top, img_w, img_h, "Sem imagem", font_size=layout_config['SKU_FONT_SIZE'])
    
    # Adicionando a caixa de texto do nome do produto
    add_textbox(slide, left + padding, top + box_h - titulo_h - sku_h - brand_h, box_w - 2 * padding, titulo_h, nome, font_size=layout_config['PRODUCT_NAME_FONT_SIZE'], bold=True, align=PP_ALIGN.LEFT)
    
    # Adicionando a caixa de texto do SKU
    sku_box = add_textbox(slide, left + padding, top + box_h - sku_h - brand_h, box_w - 2 * padding, sku_h, f"SKU: {sku}", font_size=layout_config['SKU_FONT_SIZE'], align=PP_ALIGN.LEFT)
    for p in sku_box.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = RGBColor(90, 90, 90)

    # Adicionando a nova caixa de texto da marca, com alinhamento à esquerda
    if brand:
        brand_box = add_textbox(slide, left + padding, top + box_h - brand_h, box_w - 2 * padding, brand_h, brand, font_size=layout_config['BRAND_FONT_SIZE'], bold=True, align=PP_ALIGN.LEFT)
    
    # Adicionando a caixa de texto do nome do produto
    add_textbox(slide, left + padding, top + box_h - titulo_h - sku_h - brand_h, box_w - 2 * padding, titulo_h, nome, font_size=layout_config['PRODUCT_NAME_FONT_SIZE'], bold=True, align=PP_ALIGN.LEFT)
    
    # Adicionando a caixa de texto do SKU
    sku_box = add_textbox(slide, left + padding, top + box_h - sku_h - brand_h, box_w - 2 * padding, sku_h, f"SKU: {sku}", font_size=layout_config['SKU_FONT_SIZE'], align=PP_ALIGN.LEFT)
    for p in sku_box.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = RGBColor(90, 90, 90)

    # Adicionando a nova caixa de texto da marca, com alinhamento à esquerda
    if brand:
        brand_box = add_textbox(slide, left + padding, top + box_h - brand_h, box_w - 2 * padding, brand_h, brand, font_size=layout_config['BRAND_FONT_SIZE'], bold=True, align=PP_ALIGN.LEFT)

def add_category_title(slide, title_text, layout_config):
    left, top = layout_config['MARGEM'], layout_config['TOPO'] - Inches(0.8)
    box = slide.shapes.add_textbox(left, top, Inches(10.0), Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(title_text)
    font = run.font
    font.size = layout_config['TITLE_FONT_SIZE']
    font.bold = True
    font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

def build_catalog(args, layout_config):
    logger.info("Iniciando geração do catálogo...")
    logger.debug(f"Argumentos: {vars(args)}")

    # 1) Carrega dados da planilha principal (produtos)
    df_prod, c_sku, c_nome, c_cat, c_marca_prod = load_product_data(args.excel)
    logger.info(f"Produtos carregados: {len(df_prod)} linhas válidas.")

    # 2) Carrega base de links de imagem
    df_links = load_image_links(args.img_excel)
    logger.info(f"Base de imagens carregada: {len(df_links)} SKUs com URL.")

    # 3) LEFT JOIN
    # Usa o SKU da tabela de produtos como a chave de união principal
    df = df_prod.merge(df_links, how="left", left_on=c_sku, right_on="SKU_img")
    logger.info(f"Após LEFT JOIN, continuamos com {len(df)} produtos.")

    # 4) Download de imagens
    if not args.skip_download:
        candidatos = df[df["ImageURL"].notna()].copy()
        logger.info(f"SKUs com URL disponível para download: {len(candidatos)}")

        baixados, pulados, falhas = 0, 0, 0
        for _, row in candidatos.iterrows():
            sku = str(row[c_sku]).strip()
            url = str(row["ImageURL"]).strip() if row["ImageURL"] else ""
            if not sku or not url:
                pulados += 1
                continue
            
            if find_image_for_sku(sku, args.imagens):
                pulados += 1
                logger.debug(f"Imagem já existe localmente para SKU {sku}, pulando download.")
                continue

            logger.info(f"Baixando imagem para SKU {sku}...")
            saved = download_image(url, sku, args.imagens)
            if saved:
                baixados += 1
            else:
                falhas += 1
        
        logger.info(f"Downloads concluídos. Sucesso: {baixados} | Pulados: {pulados} | Falhas: {falhas}")
    else:
        logger.info("Opção --skip-download ativada, pulando a etapa de download de imagens.")

    # 5) Prepara PowerPoint
    prs = Presentation(args.template)
    logger.info("Template PPTX carregado.")
    blank_layout = prs.slide_layouts[6]

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    GRID_COLS = layout_config['GRID_COLS']
    GRID_ROWS = layout_config['GRID_ROWS']
    MARGEM = layout_config['MARGEM']
    TOPO = layout_config['TOPO']
    BASE = layout_config['BASE']
    HGAP = layout_config['HGAP']
    VGAP = layout_config['VGAP']

    grid_w = slide_w - 2 * MARGEM
    grid_h = slide_h - TOPO - BASE
    box_w = (grid_w - (GRID_COLS - 1) * HGAP) / GRID_COLS
    box_h = (grid_h - (GRID_ROWS - 1) * VGAP) / GRID_ROWS

    # 6) Ordena e gera slides por categoria
    df = df.sort_values(by=[c_cat, c_nome, c_sku], kind="stable")
    logger.info("Ordenação concluída por categoria/nome/SKU.")
    
    if not args.include_no_image:
        df = df[df["ImageURL"].notna()]
        logger.info(f"Filtro ativo. {len(df)} produtos com URL de imagem serão incluídos.")
    else:
        logger.info("Filtro desativado. Todos os produtos serão incluídos, mesmo sem imagem.")
    
    if args.max_products > 0 and len(df) > args.max_products:
        df = df.head(args.max_products)
        logger.info(f"Limite de {args.max_products} produtos aplicado.")

    for cat, group in df.groupby(c_cat, dropna=False):
        cat_name = "" if cat is None else str(cat)
        logger.info(f"Criando slides para categoria: '{cat_name}' ({len(group)} itens).")

        items = group.to_dict("records")
        tile_index = 0
        while tile_index < len(items):
            slide = prs.slides.add_slide(blank_layout)
            add_category_title(slide, cat_name if cat_name else "Sem Categoria", layout_config)
            for r in range(GRID_ROWS):
                for c in range(GRID_COLS):
                    if tile_index >= len(items):
                        break
                    item = items[tile_index]
                    
                    sku = str(item[c_sku]).strip()
                    nome = str(item[c_nome]).strip()
                    # A marca agora é extraída da coluna correta e sem conflitos de nome
                    marca = str(item['Marca_img']).strip() if 'Marca_img' in item and item['Marca_img'] else ''
                    
                    img_path = find_image_for_sku(sku, args.imagens)

                    if not img_path:
                        logger.warning(f"Produto {sku} ('{nome}') sem imagem local.")

                    left = MARGEM + c * (box_w + HGAP)
                    top = TOPO + r * (box_h + VGAP)
                    add_product_tile(slide, left, top, box_w, box_h, img_path, nome, sku, marca, layout_config)
                    tile_index += 1
                else:
                    continue
                break

    # 7) Salva
    prs.save(args.out)
    logger.info(f"Catálogo gerado com sucesso: {args.out}")